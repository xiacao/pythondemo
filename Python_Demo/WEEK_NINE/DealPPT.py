from pptx import Presentation

#一个Presentation     对象可以有多个幻灯片，这些幻灯片储存在slides列表里面
""" Presentation    操作PPT对象
    Slides          对幻灯片进行操作(幻灯片列表)
    Shapes          对幻灯片某一区域操作(一个幻灯片中的各种元素）
    Table           表格操作
    Text            文本
    
    命令                功能
    Presentation()    创建PPT文档
    .slide_layouts[]  返回某个幻灯片母板下的slide列表(初始可以为空)
    .slides.add_slide()增加slide
    Slide.shape.title  标题
    Slide.shape.placeholders内容
    .save()存储幻灯片"""
prs = Presentation()#生成ppt对象
bullet_slide_layout = prs.slide_layouts[1]#返回第二个幻灯片母板下的slide列表

slide = prs.slides.add_slide(bullet_slide_layout)#返回一个新添加的幻灯片，继承了从slide_layout布局
shapes = slide.shapes   #返回幻灯片的形状

title_shape = shapes.title#幻灯片的标题
body_shape = shapes.placeholders[1]#幻灯片的第一空白区域

title_shape.text = 'Adding a Bullet Slide'#标题是一个对象，它有text属性
tf = body_shape.text_frame  #空白区域有五级文字
tf.text = 'Find the bullet slide layout'#空白区域的文本

p = tf.add_paragraph()#添加段落
p.text = 'Use_TextFrame.text for first bullet'
p.level = 1               #设置文字等级

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

slide = prs.slides.add_slide(bullet_slide_layout)#用bullet_slide_layout母板添加新幻灯片(空白)

prs.save('test.pptx')#保存
