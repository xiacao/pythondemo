from pptx import Presentation
from pptx.util import Inches

#一个Presentation     对象可以有多个幻灯片，这些幻灯片储存在slides列表里面
""" Presentation    操作PPT对象
    Slides          对幻灯片进行操作(幻灯片列表)
    Shapes          对幻灯片某一区域操作(一个幻灯片中的各种元素）
    Table           表格操作
    Text            文本
    
    命令                功能
    Presentation()    创建PPT文档
    .slide_layouts[]  返回某个幻灯片母板下的slide列表(初始可以为空)
    .slides.add_slide()增加slide
    Slide.shape.title  标题
    Slide.shape.placeholders内容
    .save()存储幻灯片"""
img_path = 'python_logo.gif'

prs = Presentation()#生成ppt对象
blank_slide_layout = prs.slide_layouts[6]#返回第7个幻灯片母板下的slide列表

slide = prs.slides.add_slide(blank_slide_layout)#返回一个新添加的幻灯片，继承了从slide_layout布局

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top, height=height)




slide = prs.slides.add_slide(blank_slide_layout)#用blank_slide_layout母板添加新幻灯片(空白)

prs.save('test2.pptx')#保存
